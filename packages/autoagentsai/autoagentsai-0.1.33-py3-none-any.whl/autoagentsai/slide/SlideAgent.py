from ..client import ChatClient
from ..utils.extractor import extract_json
from ..utils.convertor import convert_csv_to_json_list
from .utils_slide import (
    download_image, 
    parse_markdown_text, 
    apply_inline_formatting, 
    enable_bullet, 
    fill_existing_table, 
    find_nearest_table,
    get_value_by_path
)
import os
from typing import List
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


class SlideAgent:
    def __init__(self):
        pass

    def outline(self, prompt: str, file_path_list: List[str]):
        chat_client = ChatClient(
            agent_id="045c418f0dcf4adbb2f15031f06694d1",
            personal_auth_key="48cf18e0e0ca4b51bbf8fa60193ffb5c",
            personal_auth_secret="HWlQXZ5vxgrXDGEtTGGdsTFhJfr9rCmD",
            base_url="https://uat.agentspro.cn"
        )
    
        print(f"Debug: 准备处理 {len(file_path_list)} 个文件: {file_path_list}")
        
        content = ""
        try:
            for event in chat_client.invoke(prompt, files=file_path_list):
                if event['type'] == 'start_bubble':
                    print(f"\n{'=' * 20} 消息气泡{event['bubble_id']}开始 {'=' * 20}")
                elif event['type'] == 'token':
                    print(event['content'], end='', flush=True)
                    content += event['content']
                elif event['type'] == 'end_bubble':
                    print(f"\n{'=' * 20} 消息气泡结束 {'=' * 20}")
                elif event['type'] == 'finish':
                    print(f"\n{'=' * 20} 对话完成 {'=' * 20}")
                    break
                elif event['type'] == 'error':
                    print(f"\nDebug: 收到错误事件: {event}")
                    break
                    
        except Exception as e:
            print(f"\nDebug: ChatClient.invoke 发生异常: {type(e).__name__}: {e}")
            # 如果流出现问题，返回错误信息而不是空字符串
            if not content.strip():
                content = f"Stream error: {str(e)}"
        
        print(f"\nDebug: 最终返回内容长度: {len(content)}")
        content = extract_json(content)
        return content


    def fill(self, data: dict, template_file_path: str, output_file_path: str):
        # 加载 PPTX 模板
        prs = Presentation(template_file_path)

        # 用于存储需要清理的临时文件
        temp_files = []

        # 处理远程图片下载
        processed_data = {}
        # 支持的图片文件后缀
        image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.svg')
        
        def process_value(value):
            """递归处理数据值，支持CSV文件读取"""
            if isinstance(value, str):
                # 检查是否是CSV文件路径
                if value.endswith('.csv') and os.path.exists(value):
                    print(f"检测到CSV文件: {value}")
                    return convert_csv_to_json_list(value)
                # 检查是否是远程图片URL
                elif value.startswith(('http://', 'https://')):
                    # 检查URL是否以图片文件后缀结尾（忽略查询参数）
                    url_path = value.split('?')[0].lower()  # 去掉查询参数并转为小写
                    if url_path.endswith(image_extensions):
                        # 尝试下载远程图片
                        local_image_path = download_image(value)
                        if local_image_path:
                            temp_files.append(local_image_path)
                            print(f"成功下载图片: {value} -> {local_image_path}")
                            return local_image_path
                        else:
                            print(f"跳过下载失败的图片: {value}")
                            return None
                    else:
                        print(f"跳过非图片URL: {value} (不支持的文件类型)")
                        return value
                else:
                    return value
            elif isinstance(value, list):
                # 递归处理列表中的每个元素
                return [process_value(item) for item in value]
            elif isinstance(value, dict):
                # 递归处理字典中的每个值
                return {k: process_value(v) for k, v in value.items()}
            else:
                return value

        for key, value in data.items():
            processed_value = process_value(value)
            if processed_value is not None:
                processed_data[key] = processed_value

        # 1. 表格填充
        table_requests = []  # [(占位符形状, key, data)]
        all_tables = []  # 所有表格形状
        
        print(f"开始扫描PPT模板中的占位符...")
        for slide_idx, slide in enumerate(prs.slides):
            print(f"扫描第 {slide_idx + 1} 页...")
            for shape in slide.shapes:
                # 收集表格占位符
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text.startswith("{{") and text.endswith("}}"):
                        print(f"  找到占位符: {text}")
                    if text.startswith("{{#") and text.endswith("}}"):
                        path = text[3:-2].strip()  # 去掉 {{# 和 }}
                        print(f"找到表格占位符: {{#{path}}}")
                        table_data = get_value_by_path(processed_data, path)
                        if table_data is not None:
                            print(f"表格占位符 {{#{path}}} 数据解析成功，{len(table_data) if isinstance(table_data, list) else 1} 条记录")
                            table_requests.append((shape, path, table_data))
                        else:
                            print(f"表格占位符 {{#{path}}} 数据解析失败")
                
                # 收集所有表格
                if shape.has_table:
                    all_tables.append(shape)
        
        # 为每个表格占位符找到最近的表格并填充
        shapes_to_remove = []
        processed_tables = set()
        
        for placeholder_shape, path, table_data in table_requests:
            available_tables = [t for t in all_tables if id(t) not in processed_tables]
            if not available_tables:
                available_tables = all_tables
            
            nearest_table_shape = find_nearest_table(placeholder_shape, available_tables)
            if nearest_table_shape:
                print(f"占位符 '{{#{path}}}' 匹配到最近的表格")
                fill_existing_table(nearest_table_shape.table, table_data)
                processed_tables.add(id(nearest_table_shape))
            
            shapes_to_remove.append(placeholder_shape)
        
        # 删除表格占位符文本框
        for shape in shapes_to_remove:
            shape._element.getparent().remove(shape._element)
        
        # 2. 文本、图片填充
        for slide in prs.slides:
            for shape in list(slide.shapes):  # list() to allow removal
                if not shape.has_text_frame:
                    continue
            
                text = shape.text.strip()
                if text.startswith("{{") and text.endswith("}}"):
                    path = text[2:-2].strip()  # 去掉 {{}}
                    content_type = "text"

                    # 判断类型前缀
                    if path.startswith("@"):
                        path = path[1:]
                        content_type = "image"
                    elif path.startswith("#"):
                        # 表格已经在上面处理过了，跳过
                        continue

                    value = get_value_by_path(processed_data, path)
                    if value is None:
                        continue

                    if content_type == "text":
                        # 检查是否包含Markdown格式
                        if isinstance(value, str) and any(marker in value for marker in ['*', '#', '`', '\n']):
                            # 使用Markdown解析
                            parse_markdown_text(shape.text_frame, value)
                        elif isinstance(value, list):
                            # 处理列表数据，每项作为bullet point
                            tf = shape.text_frame
                            tf.clear()
                            for i, item in enumerate(value):
                                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                                if isinstance(item, str) and any(marker in item for marker in ['*', '#', '`']):
                                    apply_inline_formatting(p, item)
                                else:
                                    p.text = str(item)
                                    p.font.size = Pt(14)
                                    p.alignment = PP_ALIGN.LEFT
                                    enable_bullet(p)
                        else:
                            # 普通文本
                            shape.text_frame.text = str(value)

                    elif content_type == "image":
                        # 获取位置并删除原文本框
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        slide.shapes._spTree.remove(shape._element)
                            
                        # 确保图片路径存在
                        if os.path.exists(value):
                            slide.shapes.add_picture(value, left, top, width=width, height=height)
                            print(f"成功替换图片: {path}")
                        else:
                            print(f"警告: 图片文件不存在: {value}")

        # 保存为新PPT
        prs.save(output_file_path)
        
        # 清理临时文件
        for temp_file in temp_files:
            try:
                os.unlink(temp_file)
                print(f"清理临时文件: {temp_file}")
            except Exception as e:
                print(f"清理临时文件失败: {temp_file}, 错误: {e}")