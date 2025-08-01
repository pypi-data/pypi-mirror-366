from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from pptx.oxml import parse_xml
import re
import tempfile
import os
import shutil
import requests
from pathlib import Path
from typing import Optional, Any


def parse_markdown_text(text_frame, markdown_text, font_size=14):
    """
    解析Markdown文本并应用到PowerPoint文本框
    支持：
    - * bullet points
    - **粗体**
    - *斜体*
    - `代码`
    - # 标题
    """
    text_frame.clear()
    
    lines = markdown_text.split('\n')
    
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
            
        # 添加段落
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # 处理标题 (# ## ###)
        if line.startswith('#'):
            level = 0
            while level < len(line) and line[level] == '#':
                level += 1
            title_text = line[level:].strip()
            p.text = title_text
            p.font.size = Pt(font_size + (4 - level) * 2)  # 标题字体更大
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
            continue
        
        # 处理bullet points
        if line.startswith('* ') or line.startswith('- '):
            bullet_text = line[2:].strip()
            p.text = bullet_text
            p.font.size = Pt(font_size)
            p.alignment = PP_ALIGN.LEFT
            enable_bullet(p)
            
            # 处理bullet point内的格式
            apply_inline_formatting(p, bullet_text)
            continue
        
        # 处理普通文本
        p.text = line
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.LEFT
        
        # 处理内联格式
        apply_inline_formatting(p, line)


def apply_inline_formatting(paragraph, text):
    """
    应用内联格式：粗体、斜体、代码
    """
    # 清空段落文本，重新构建带格式的文本
    paragraph.clear()
    
    # 使用正则表达式解析格式
    # 匹配 **粗体**、*斜体*、`代码`
    pattern = r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)'
    
    parts = re.split(pattern, text)
    
    for part in parts:
        if not part:
            continue
            
        run = paragraph.add_run()
        
        if part.startswith('**') and part.endswith('**'):
            # 粗体
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith('*') and part.endswith('*'):
            # 斜体
            run.text = part[1:-1]
            run.font.italic = True
        elif part.startswith('`') and part.endswith('`'):
            # 代码
            run.text = part[1:-1]
            run.font.name = 'Consolas'
            run.font.color.rgb = RGBColor(220, 20, 60)  # 深红色
        else:
            # 普通文本
            run.text = part


def enable_bullet(paragraph, bullet_char="•"):
    """为段落启用项目符号"""
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    buChar = parse_xml(f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{bullet_char}"/>')
    pPr.insert(0, buChar)


def fill_existing_table(table, data, font_size=12):
    """
    将 data 填充到 pptx 表格中，第二行作为模板行。
    支持字段替换：[name], [count], [=count*price], [@picture]
    """
    from pptx.util import Pt
    import re

    def eval_formula(expr, context):
        try:
            return str(eval(expr, {}, context))
        except:
            return expr

    # 查找包含占位符的行作为模板行
    template_row_idx = None
    for i, row in enumerate(table.rows):
        for cell in row.cells:
            if '[' in cell.text and ']' in cell.text:
                template_row_idx = i
                break
        if template_row_idx is not None:
            break
    
    if template_row_idx is None:
        return
        
    template_row = table.rows[template_row_idx]
    col_templates = [cell.text for cell in template_row.cells]

    for row_data in data:
        # 使用底层XML操作添加新行
        tbl = table._tbl
        new_tr = parse_xml(template_row._tr.xml)
        tbl.append(new_tr)
        
        # 获取新添加的行
        new_row_idx = len(table.rows) - 1
        row = table.rows[new_row_idx]
        
        for j, tmpl in enumerate(col_templates):
            text = tmpl
            
            # 处理图片占位符 [@picture]
            if re.search(r'\[@\w+\]', text):
                # 找到图片占位符
                img_match = re.search(r'\[@(\w+)\]', text)
                if img_match:
                    img_key = img_match.group(1)
                    if img_key in row_data:
                        # 清空单元格文本
                        row.cells[j].text = ""
                        # 这里可以添加图片到单元格的逻辑
                        # 目前先显示图片文件名
                        row.cells[j].text = f"图片: {row_data[img_key]}"
                    continue
            
            # 字段替换
            for key, val in row_data.items():
                text = text.replace(f"[{key}]", str(val))
            
            # 表达式处理（如 [=count*price]）
            match = re.findall(r"\[=([^\]]+)\]", text)
            for m in match:
                result = eval_formula(m, row_data)
                text = text.replace(f"[={m}]", result)
            
            # 检查是否包含Markdown格式
            if any(marker in text for marker in ['*', '#', '`']):
                # 使用Markdown解析
                parse_markdown_text(row.cells[j].text_frame, text, font_size)
            else:
                # 普通文本
                row.cells[j].text = text
                para = row.cells[j].text_frame.paragraphs[0]
                para.font.size = Pt(font_size)
                para.alignment = PP_ALIGN.CENTER
                row.cells[j].text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 删除模板行
    tbl.remove(template_row._tr)


def find_nearest_table(placeholder_shape, all_tables):
    """
    根据最近距离原则找到对应的表格
    """
    if not all_tables:
        return None
    
    placeholder_pos = (placeholder_shape.left, placeholder_shape.top)
    
    def calculate_distance(table_shape):
        table_pos = (table_shape.left, table_shape.top)
        return ((placeholder_pos[0] - table_pos[0]) ** 2 + 
                (placeholder_pos[1] - table_pos[1]) ** 2) ** 0.5
    
    return min(all_tables, key=calculate_distance)


def create_temp_file(file_content, file_extension='.pptx'):
    """
    创建临时文件并写入内容
    
    Args:
        file_content: 文件内容（字节或字符串）
        file_extension: 文件扩展名，默认为.pptx
    
    Returns:
        str: 临时文件路径
    """
    temp_fd, temp_path = tempfile.mkstemp(suffix=file_extension)
    try:
        with os.fdopen(temp_fd, 'wb' if isinstance(file_content, bytes) else 'w') as tmp_file:
            tmp_file.write(file_content)
        return temp_path
    except Exception as e:
        os.unlink(temp_path)
        raise e


def copy_file_to_temp(source_path, file_extension=None):
    """
    将文件复制到临时目录
    
    Args:
        source_path: 源文件路径
        file_extension: 目标文件扩展名，默认保持原扩展名
    
    Returns:
        str: 临时文件路径
    """
    if file_extension is None:
        file_extension = Path(source_path).suffix
    
    temp_fd, temp_path = tempfile.mkstemp(suffix=file_extension)
    os.close(temp_fd)
    
    try:
        shutil.copy2(source_path, temp_path)
        return temp_path
    except Exception as e:
        os.unlink(temp_path)
        raise e


def cleanup_temp_file(temp_path):
    """
    清理临时文件
    
    Args:
        temp_path: 临时文件路径
    """
    try:
        if os.path.exists(temp_path):
            os.unlink(temp_path)
    except Exception:
        pass  # 忽略清理错误


def fill_pptx(template_path, output_path, replacements):
    """
    使用替换数据填充PowerPoint模板
    
    Args:
        template_path: 模板文件路径
        output_path: 输出文件路径
        replacements: 替换数据字典
    """
    prs = Presentation(template_path)

    # 1.表格填充
    # 第一步：收集所有表格占位符和表格信息
    table_requests = []  # [(占位符形状, key, data)]
    all_tables = []  # 所有表格形状
    
    for slide in prs.slides:
        for shape in slide.shapes:
            # 收集表格占位符
            if shape.has_text_frame:
                text = shape.text.strip()
                if text.startswith("{{#") and text.endswith("}}"):
                    key = text[3:-2].strip()  # 去掉 {{# 和 }}
                    if key in replacements:
                        table_requests.append((shape, key, replacements[key]))
            
            # 收集所有表格
            if shape.has_table:
                all_tables.append(shape)
    
    # 第二步：为每个表格占位符找到最近的表格并填充
    shapes_to_remove = []
    processed_tables = set()  # 使用shape的id来跟踪已处理的表格
    
    for placeholder_shape, key, data in table_requests:
        # 找到最近的未处理表格
        available_tables = [t for t in all_tables if id(t) not in processed_tables]
        if not available_tables:
            # 如果所有表格都被处理过，则使用最近的表格（允许重复使用）
            available_tables = all_tables
        
        nearest_table_shape = find_nearest_table(placeholder_shape, available_tables)
        if nearest_table_shape:
            print(f"占位符 '{{#{key}}}' 匹配到最近的表格")
            fill_existing_table(nearest_table_shape.table, data)
            processed_tables.add(id(nearest_table_shape))
        
        shapes_to_remove.append(placeholder_shape)
    
    # 第三步：删除表格占位符文本框
    for shape in shapes_to_remove:
        shape._element.getparent().remove(shape._element)
    
    # 2.文本、图片填充
    for slide in prs.slides:
        for shape in list(slide.shapes):  # list() to allow removal
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if text.startswith("{{") and text.endswith("}}"):
                key = text[2:-2].strip()  # 去掉 {{}}
                content_type = "text"

                # 判断类型前缀
                if key.startswith("@"):
                    key = key[1:]
                    content_type = "image"
                elif key.startswith("#"):
                    # 表格已经在上面处理过了，跳过
                    continue

                value = replacements.get(key)
                if value is None:
                    continue

                if content_type == "text":
                    # 检查是否包含Markdown格式
                    if isinstance(value, str) and any(marker in value for marker in ['*', '#', '`', '\n']):
                        # 使用Markdown解析
                        parse_markdown_text(shape.text_frame, value)
                    elif isinstance(value, list):
                        # 处理列表数据，每项作为bullet point
                        tf = shape.text_frame
                        tf.clear()
                        for i, item in enumerate(value):
                            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                            if isinstance(item, str) and any(marker in item for marker in ['*', '#', '`']):
                                # item包含Markdown，解析格式
                                apply_inline_formatting(p, item)
                            else:
                                p.text = str(item)
                                p.font.size = Pt(14)
                                p.alignment = PP_ALIGN.LEFT
                                enable_bullet(p)
                    else:
                        # 普通文本
                        shape.text_frame.text = str(value)

                elif content_type == "image":
                    # 获取位置并删除原文本框
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    slide.shapes._spTree.remove(shape._element)
                    slide.shapes.add_picture(value, left, top, width=width, height=height)

    prs.save(output_path)


def fill_pptx_with_temp_files(template_file_content, replacements, output_path=None):
    """
    使用临时文件处理PowerPoint模板填充
    
    Args:
        template_file_content: 模板文件内容（字节）
        replacements: 替换数据字典
        output_path: 输出文件路径，如果为None则返回临时文件路径
    
    Returns:
        str: 输出文件路径
    """
    # 创建临时模板文件
    temp_template_path = create_temp_file(template_file_content, '.pptx')
    
    try:
        # 如果没有指定输出路径，创建临时输出文件
        if output_path is None:
            temp_fd, output_path = tempfile.mkstemp(suffix='.pptx')
            os.close(temp_fd)
        
        # 执行填充
        fill_pptx(temp_template_path, output_path, replacements)
        
        return output_path
        
    finally:
        # 清理临时模板文件
        cleanup_temp_file(temp_template_path)


def download_image(url: str) -> Optional[str]:
    """
    下载远程图片到临时文件
    
    Args:
        url: 图片URL
    
    Returns:
        str: 临时文件路径，如果下载失败返回None
    """
    
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()
        
        # 验证content-type是否为图片
        content_type = response.headers.get('content-type', '').lower()
        valid_image_types = [
            'image/jpeg', 'image/jpg', 'image/png', 'image/gif', 
            'image/bmp', 'image/webp', 'image/tiff', 'image/svg+xml'
        ]
        
        if not any(img_type in content_type for img_type in valid_image_types):
            print(f"跳过非图片内容: {url}, content-type: {content_type}")
            return None
        
        # 创建临时文件
        suffix = '.jpg'  # 默认后缀
        if 'image/png' in content_type:
            suffix = '.png'
        elif 'image/gif' in content_type:
            suffix = '.gif'
        elif 'image/webp' in content_type:
            suffix = '.webp'
        elif 'image/bmp' in content_type:
            suffix = '.bmp'
        elif 'image/tiff' in content_type:
            suffix = '.tiff'
        elif 'image/svg' in content_type:
            suffix = '.svg'
        
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        
        # 下载图片内容
        for chunk in response.iter_content(chunk_size=8192):
            temp_file.write(chunk)
        
        temp_file.close()
        return temp_file.name
        
    except Exception as e:
        print(f"下载图片失败: {url}, 错误: {e}")
        return None

def get_value_by_path(data: dict, path: str) -> Any:
    """
    根据路径表达式从数据中获取值
    支持格式：page[0].title, page[1].sections[0].content 等
    """
    try:
        # 将路径分解为步骤
        current = data
        
        # 使用正则表达式匹配路径中的各个部分
        # 匹配形如 "key" 或 "key[index]" 的模式
        pattern = r'([a-zA-Z_][a-zA-Z0-9_]*)?(?:\[(\d+)\])?'
        steps = re.findall(r'([a-zA-Z_][a-zA-Z0-9_]*(?:\[\d+\])?)', path)
        
        for step in steps:
            # 检查是否包含数组索引
            if '[' in step and ']' in step:
                key_part = step.split('[')[0]
                index_part = step.split('[')[1].rstrip(']')
                index = int(index_part)
                
                if key_part:
                    current = current[key_part]
                current = current[index]
            else:
                current = current[step]
        
        return current
    except (KeyError, IndexError, ValueError, TypeError) as e:
        print(f"路径解析错误: {path}, 错误: {e}")
        return None

