"""PowerPoint file reader implementation."""

from pathlib import Path
from typing import List, Optional, Any
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from vector_db_query.document_processor.office_readers import PresentationReader
from vector_db_query.document_processor.exceptions import DocumentReadError
from vector_db_query.utils.logger import get_logger

logger = get_logger(__name__)


class PowerPointReader(PresentationReader):
    """Reader for PowerPoint files (.pptx, .ppt)."""
    
    def can_read(self, file_path: Path) -> bool:
        """Check if this reader can handle the file type."""
        return file_path.suffix.lower() in self.supported_extensions
        
    def read(self, file_path: Path) -> str:
        """Read PowerPoint file and extract text content."""
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.pptx':
                return self._read_pptx(file_path)
            elif extension == '.ppt':
                # For legacy .ppt files, try to read as .pptx first
                # python-pptx might handle some .ppt files
                try:
                    return self._read_pptx(file_path)
                except Exception:
                    # If that fails, try legacy method
                    return self._read_ppt_legacy(file_path)
            else:
                raise DocumentReadError(
                    f"Unsupported PowerPoint format: {extension}",
                    file_path=str(file_path)
                )
        except Exception as e:
            raise DocumentReadError(
                f"Failed to read PowerPoint file: {e}",
                file_path=str(file_path)
            )
            
    @property
    def supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return ['.pptx', '.ppt']
        
    def _read_pptx(self, file_path: Path) -> str:
        """Read .pptx file using python-pptx."""
        logger.info(f"Reading PowerPoint file: {file_path.name}")
        
        prs = Presentation(str(file_path))
        text_parts = []
        
        # Extract metadata
        self._metadata = self._extract_metadata(file_path)
        self._metadata.update({
            'total_slides': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height
        })
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = self._extract_slide_content(slide, slide_num)
            if slide_parts:
                text_parts.append(self.slide_separator.join(slide_parts))
                
        return self.slide_separator.join(text_parts)
        
    def _extract_slide_content(self, slide, slide_num: int) -> List[str]:
        """Extract content from a single slide."""
        parts = []
        
        # Get slide title
        title = self._get_slide_title(slide)
        
        # Get all text from shapes
        content_parts = []
        for shape in slide.shapes:
            shape_text = self._extract_shape_text(shape)
            if shape_text and shape_text != title:  # Avoid duplicating title
                content_parts.append(shape_text)
                
        # Get notes if enabled
        notes = None
        if self.include_notes and slide.has_notes_slide:
            notes = self._extract_notes(slide.notes_slide)
            
        # Format the slide
        slide_text = self._format_slide(
            slide_num=slide_num,
            title=title or "",
            content="\n".join(content_parts),
            notes=notes
        )
        
        if slide_text.strip():
            parts.append(slide_text)
            
        return parts
        
    def _get_slide_title(self, slide) -> Optional[str]:
        """Extract slide title if present."""
        if slide.shapes.title:
            return slide.shapes.title.text
        return None
        
    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape."""
        text_parts = []
        
        # Handle text frames
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = ''.join(run.text for run in paragraph.runs)
                if para_text.strip():
                    text_parts.append(para_text.strip())
                    
        # Handle tables
        elif shape.has_table:
            table_text = self._extract_table_text(shape.table)
            if table_text:
                text_parts.append(table_text)
                
        # Handle grouped shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                sub_text = self._extract_shape_text(sub_shape)
                if sub_text:
                    text_parts.append(sub_text)
                    
        return '\n'.join(text_parts)
        
    def _extract_table_text(self, table) -> str:
        """Extract text from a table."""
        rows_text = []
        
        for row in table.rows:
            cells_text = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    cells_text.append(cell_text)
            if cells_text:
                rows_text.append(' | '.join(cells_text))
                
        return '\n'.join(rows_text)
        
    def _extract_notes(self, notes_slide) -> Optional[str]:
        """Extract speaker notes from notes slide."""
        notes_parts = []
        
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ''.join(run.text for run in paragraph.runs)
                    if para_text.strip():
                        notes_parts.append(para_text.strip())
                        
        return '\n'.join(notes_parts) if notes_parts else None
        
    def _read_ppt_legacy(self, file_path: Path) -> str:
        """Attempt to read legacy .ppt files.
        
        This is a fallback method that tries to extract any readable text.
        """
        logger.warning(f"Using legacy method for .ppt file: {file_path.name}")
        
        # For now, we'll return a message indicating limited support
        # In a production system, you might use additional libraries
        # like python-ppt or extract text using system tools
        
        return (
            f"Legacy PowerPoint file: {file_path.name}\n"
            "Note: Full support for .ppt files requires additional tools.\n"
            "Please convert to .pptx format for best results."
        )
        
    def _process_content(self, content: Any) -> str:
        """Process the raw content into text.
        
        This is implemented through the read methods above.
        """
        # Not used in this implementation
        pass