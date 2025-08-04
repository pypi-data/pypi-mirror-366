"""Create a test PowerPoint file."""

from pptx import Presentation
from pptx.util import Inches

# Create presentation
prs = Presentation()

# Slide 1: Title slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Test Presentation"
subtitle.text = "Created for Vector DB Query System"

# Slide 2: Content slide
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide2.shapes.title
content = slide2.placeholders[1]
title.text = "Key Features"
content.text = "• Document Processing\n• Vector Search\n• LLM Integration"

# Add speaker notes
notes_slide = slide2.notes_slide
notes_slide.notes_text_frame.text = "These are speaker notes for the features slide."

# Save
prs.save('tests/sample_files/test_presentation.pptx')
print("Test PowerPoint created: tests/sample_files/test_presentation.pptx")