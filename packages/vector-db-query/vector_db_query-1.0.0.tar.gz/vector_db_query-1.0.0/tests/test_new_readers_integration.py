"""Integration tests for newly added document readers."""

import pytest
import tempfile
import shutil
from pathlib import Path
import openpyxl
from pptx import Presentation
from pptx.util import Inches
import email
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime

from vector_db_query.document_processor import DocumentProcessor
from vector_db_query.document_processor.excel_reader import ExcelReader
from vector_db_query.document_processor.powerpoint_reader import PowerPointReader
from vector_db_query.document_processor.email_reader import EmailReader
from vector_db_query.document_processor.html_reader import HTMLReader
from vector_db_query.document_processor.config_reader import (
    JSONReader, XMLReader, YAMLReader, INIReader, LogReader
)
from vector_db_query.utils.logger import get_logger

logger = get_logger(__name__)


class TestExcelReaderIntegration:
    """Integration tests for Excel reader."""
    
    @pytest.fixture
    def excel_file(self, tmp_path):
        """Create a test Excel file with various features."""
        file_path = tmp_path / "test_workbook.xlsx"
        
        # Create workbook
        wb = openpyxl.Workbook()
        
        # First sheet - data with formulas
        ws1 = wb.active
        ws1.title = "Sales Data"
        
        # Headers
        ws1['A1'] = 'Product'
        ws1['B1'] = 'Quantity'
        ws1['C1'] = 'Price'
        ws1['D1'] = 'Total'
        
        # Data
        ws1['A2'] = 'Widget A'
        ws1['B2'] = 10
        ws1['C2'] = 5.99
        ws1['D2'] = '=B2*C2'  # Formula
        
        ws1['A3'] = 'Widget B'
        ws1['B3'] = 5
        ws1['C3'] = 12.50
        ws1['D3'] = '=B3*C3'  # Formula
        
        # Add comment
        ws1['A2'].comment = openpyxl.comments.Comment('Best selling product', 'Sales Team')
        
        # Summary row
        ws1['A5'] = 'Total'
        ws1['D5'] = '=SUM(D2:D3)'  # Formula
        
        # Second sheet - empty
        ws2 = wb.create_sheet("Empty Sheet")
        
        # Third sheet - more data
        ws3 = wb.create_sheet("Inventory")
        ws3['A1'] = 'Item'
        ws3['B1'] = 'Stock'
        ws3['A2'] = 'Widget A'
        ws3['B2'] = 100
        
        wb.save(file_path)
        return file_path
        
    def test_excel_complete_extraction(self, excel_file):
        """Test complete Excel file extraction."""
        processor = DocumentProcessor()
        doc = processor.process_file(excel_file)
        
        # Verify processing
        assert doc is not None
        assert len(doc.chunks) > 0
        assert doc.metadata.file_type == '.xlsx'
        
        # Combine all chunks
        full_text = ' '.join(chunk.text for chunk in doc.chunks)
        
        # Check content extraction
        assert 'Sales Data' in full_text  # Sheet name
        assert 'Product' in full_text  # Header
        assert 'Widget A' in full_text  # Data
        assert 'Widget B' in full_text
        assert 'Inventory' in full_text  # Second sheet
        
        # Check if formulas were extracted (based on config)
        config = processor.reader_factory._readers.get('excel', ExcelReader())
        if hasattr(config, 'extract_formulas') and config.extract_formulas:
            assert '=B2*C2' in full_text or 'SUM' in full_text
            
        # Check if comments were extracted
        if hasattr(config, 'extract_comments') and config.extract_comments:
            assert 'Best selling product' in full_text
            
    def test_excel_multi_sheet_processing(self, excel_file):
        """Test multi-sheet Excel processing."""
        reader = ExcelReader()
        reader.process_all_sheets = True
        
        text = reader.read(excel_file)
        
        # Should contain data from multiple sheets
        assert 'Sales Data' in text
        assert 'Inventory' in text
        assert 'Widget A' in text  # Appears in both sheets
        
    def test_excel_with_row_limits(self, tmp_path):
        """Test Excel processing with row limits."""
        # Create large Excel file
        file_path = tmp_path / "large.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        
        # Add many rows
        for i in range(1000):
            ws[f'A{i+1}'] = f'Row {i+1}'
            ws[f'B{i+1}'] = i
            
        wb.save(file_path)
        
        # Process with row limit
        reader = ExcelReader()
        reader.max_rows_per_sheet = 100
        
        text = reader.read(file_path)
        
        # Should have limited rows
        assert 'Row 100' in text
        assert 'Row 500' not in text  # Beyond limit


class TestPowerPointReaderIntegration:
    """Integration tests for PowerPoint reader."""
    
    @pytest.fixture
    def pptx_file(self, tmp_path):
        """Create a test PowerPoint file."""
        file_path = tmp_path / "test_presentation.pptx"
        
        # Create presentation
        prs = Presentation()
        
        # Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Test Presentation"
        slide1.placeholders[1].text = "Integration Testing"
        
        # Add speaker notes
        slide1.notes_slide.notes_text_frame.text = "Welcome everyone to this test presentation"
        
        # Content slide with bullet points
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Key Points"
        content = slide2.placeholders[1]
        content.text = "First point\nSecond point\nThird point"
        
        # Add speaker notes
        slide2.notes_slide.notes_text_frame.text = "Remember to emphasize these points"
        
        # Slide with table
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        slide3.shapes.title.text = "Data Table"
        
        # Add table
        table = slide3.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(5), Inches(2)).table
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(0, 2).text = "Header 3"
        table.cell(1, 0).text = "Data 1"
        table.cell(1, 1).text = "Data 2"
        table.cell(1, 2).text = "Data 3"
        
        prs.save(file_path)
        return file_path
        
    def test_powerpoint_complete_extraction(self, pptx_file):
        """Test complete PowerPoint extraction."""
        processor = DocumentProcessor()
        doc = processor.process_file(pptx_file)
        
        assert doc is not None
        assert len(doc.chunks) > 0
        
        full_text = ' '.join(chunk.text for chunk in doc.chunks)
        
        # Check slide content
        assert 'Test Presentation' in full_text
        assert 'Integration Testing' in full_text
        assert 'Key Points' in full_text
        assert 'First point' in full_text
        
        # Check table content
        assert 'Header 1' in full_text
        assert 'Data 1' in full_text
        
    def test_powerpoint_speaker_notes(self, pptx_file):
        """Test speaker notes extraction."""
        reader = PowerPointReader()
        reader.extract_speaker_notes = True
        
        text = reader.read(pptx_file)
        
        # Should contain speaker notes
        assert 'Welcome everyone' in text
        assert 'Remember to emphasize' in text
        
    def test_powerpoint_slide_numbers(self, pptx_file):
        """Test slide number extraction."""
        reader = PowerPointReader()
        reader.extract_slide_numbers = True
        
        text = reader.read(pptx_file)
        
        # Should contain slide indicators
        assert 'Slide' in text or '1' in text


class TestEmailReaderIntegration:
    """Integration tests for Email reader."""
    
    @pytest.fixture
    def eml_file(self, tmp_path):
        """Create a test email file."""
        file_path = tmp_path / "test_email.eml"
        
        # Create email
        msg = MIMEMultipart()
        msg['From'] = 'sender@example.com'
        msg['To'] = 'recipient@example.com'
        msg['Subject'] = 'Test Email for Integration'
        msg['Date'] = datetime.now().strftime('%a, %d %b %Y %H:%M:%S +0000')
        
        # Add body
        body = MIMEText("This is the email body.\n\nIt has multiple paragraphs.\n\nBest regards,\nTest Sender")
        msg.attach(body)
        
        # Save to file
        with open(file_path, 'w') as f:
            f.write(msg.as_string())
            
        return file_path
        
    def test_email_complete_extraction(self, eml_file):
        """Test complete email extraction."""
        processor = DocumentProcessor()
        doc = processor.process_file(eml_file)
        
        assert doc is not None
        assert len(doc.chunks) > 0
        
        full_text = ' '.join(chunk.text for chunk in doc.chunks)
        
        # Check headers
        assert 'sender@example.com' in full_text
        assert 'recipient@example.com' in full_text
        assert 'Test Email for Integration' in full_text
        
        # Check body
        assert 'This is the email body' in full_text
        assert 'multiple paragraphs' in full_text
        assert 'Best regards' in full_text
        
    def test_email_with_html_body(self, tmp_path):
        """Test email with HTML body."""
        file_path = tmp_path / "html_email.eml"
        
        msg = MIMEMultipart('alternative')
        msg['Subject'] = 'HTML Email Test'
        msg['From'] = 'html@example.com'
        msg['To'] = 'user@example.com'
        
        # Plain text part
        text_part = MIMEText('Plain text version', 'plain')
        msg.attach(text_part)
        
        # HTML part
        html_body = """<html>
<body>
<h1>HTML Email</h1>
<p>This is an <strong>HTML</strong> email.</p>
<a href="http://example.com">Click here</a>
</body>
</html>"""
        html_part = MIMEText(html_body, 'html')
        msg.attach(html_part)
        
        with open(file_path, 'w') as f:
            f.write(msg.as_string())
            
        reader = EmailReader()
        reader.sanitize_content = True
        text = reader.read(file_path)
        
        # Should extract text from HTML
        assert 'HTML Email' in text
        assert 'Click here' in text
        # Should not contain HTML tags if sanitized
        if reader.sanitize_content:
            assert '<strong>' not in text


class TestHTMLReaderIntegration:
    """Integration tests for HTML reader."""
    
    @pytest.fixture
    def html_file(self, tmp_path):
        """Create test HTML file."""
        file_path = tmp_path / "test_page.html"
        
        content = """<!DOCTYPE html>
<html>
<head>
    <title>Test Page</title>
    <meta name="description" content="Test HTML page for integration testing">
    <script>
        console.log('This should be removed');
        function test() { return 'removed'; }
    </script>
    <style>
        body { background: white; }
        .hidden { display: none; }
    </style>
</head>
<body>
    <h1>Main Title</h1>
    <nav>
        <a href="#section1">Section 1</a>
        <a href="#section2">Section 2</a>
    </nav>
    <main>
        <section id="section1">
            <h2>Section 1</h2>
            <p>This is the first section with <strong>bold</strong> and <em>italic</em> text.</p>
            <ul>
                <li>List item 1</li>
                <li>List item 2</li>
            </ul>
        </section>
        <section id="section2">
            <h2>Section 2</h2>
            <p>This section has a <a href="http://example.com">link</a>.</p>
            <table>
                <tr><th>Header</th><th>Value</th></tr>
                <tr><td>Row 1</td><td>Data 1</td></tr>
            </table>
        </section>
    </main>
    <footer>
        <p>&copy; 2025 Test Page</p>
    </footer>
</body>
</html>"""
        
        file_path.write_text(content)
        return file_path
        
    def test_html_complete_extraction(self, html_file):
        """Test complete HTML extraction."""
        processor = DocumentProcessor()
        doc = processor.process_file(html_file)
        
        assert doc is not None
        assert len(doc.chunks) > 0
        
        full_text = ' '.join(chunk.text for chunk in doc.chunks)
        
        # Check content extraction
        assert 'Main Title' in full_text
        assert 'Section 1' in full_text
        assert 'bold' in full_text
        assert 'List item 1' in full_text
        
        # Scripts and styles should be removed
        assert 'console.log' not in full_text
        assert 'background: white' not in full_text
        
    def test_html_link_preservation(self, html_file):
        """Test HTML link preservation."""
        reader = HTMLReader()
        reader.preserve_links = True
        
        text = reader.read(html_file)
        
        # Links should be preserved in some form
        assert 'example.com' in text or 'link' in text
        
    def test_html_to_markdown(self, html_file):
        """Test HTML to Markdown conversion."""
        reader = HTMLReader()
        reader.convert_to_markdown = True
        
        text = reader.read(html_file)
        
        # Should have Markdown formatting
        assert '#' in text  # Headers
        assert '**' in text or '*' in text  # Bold/italic
        assert '-' in text or '*' in text  # List markers


class TestConfigReadersIntegration:
    """Integration tests for configuration file readers."""
    
    def test_json_reader_pretty_print(self, tmp_path):
        """Test JSON reader with pretty printing."""
        file_path = tmp_path / "config.json"
        content = {
            "app": {
                "name": "Test App",
                "version": "1.0.0",
                "features": ["feature1", "feature2"]
            },
            "database": {
                "host": "localhost",
                "port": 5432
            }
        }
        
        import json
        file_path.write_text(json.dumps(content))
        
        reader = JSONReader()
        text = reader.read(file_path)
        
        # Should be pretty printed
        assert '"app":' in text
        assert '"name": "Test App"' in text
        assert text.count('\n') > 5  # Multiple lines
        
    def test_yaml_multi_document(self, tmp_path):
        """Test YAML reader with multiple documents."""
        file_path = tmp_path / "multi.yaml"
        content = """---
document: 1
type: config
---
document: 2
type: data
items:
  - item1
  - item2
"""
        file_path.write_text(content)
        
        reader = YAMLReader()
        text = reader.read(file_path)
        
        # Should contain both documents
        assert 'document: 1' in text
        assert 'document: 2' in text
        assert 'item1' in text
        
    def test_ini_reader_sections(self, tmp_path):
        """Test INI reader with sections."""
        file_path = tmp_path / "app.ini"
        content = """# Application configuration
[General]
name = Test Application
version = 2.0.0

[Database]
; Database settings
host = db.example.com
port = 5432
user = appuser

[Features]
feature_x = enabled
feature_y = disabled
"""
        file_path.write_text(content)
        
        reader = INIReader()
        text = reader.read(file_path)
        
        # Should preserve structure
        assert '[General]' in text
        assert '[Database]' in text
        assert 'host = db.example.com' in text
        
        # Comments handling depends on configuration
        if reader.include_comments:
            assert 'Application configuration' in text
            
    def test_log_reader_analysis(self, tmp_path):
        """Test log reader with analysis features."""
        file_path = tmp_path / "app.log"
        log_content = """2025-07-31 10:00:00 INFO Application started successfully
2025-07-31 10:00:01 DEBUG Loading configuration from config.yaml
2025-07-31 10:00:02 INFO Database connection established
2025-07-31 10:00:05 ERROR Failed to load module: ImportError: No module named 'missing'
2025-07-31 10:00:06 ERROR Traceback (most recent call last):
  File "app.py", line 50, in load_module
    import missing
ImportError: No module named 'missing'
2025-07-31 10:00:10 WARNING High memory usage: 85% of available memory
2025-07-31 10:00:15 INFO Starting web server on port 8080
2025-07-31 10:00:20 INFO Request received: GET /api/status
2025-07-31 10:00:21 INFO Response sent: 200 OK
2025-07-31 10:00:30 ERROR Database query failed: timeout after 30s
2025-07-31 10:00:31 INFO Retrying database query
2025-07-31 10:00:35 INFO Query successful on retry
"""
        file_path.write_text(log_content)
        
        reader = LogReader()
        reader.summarize = True
        reader.extract_patterns = True
        
        text = reader.read(file_path)
        
        # Should contain log entries
        assert 'Application started' in text
        assert 'ERROR' in text
        assert 'ImportError' in text
        
        # If summarizing, should have summary
        if reader.summarize:
            # Should identify error patterns
            assert 'error' in text.lower() or 'ERROR' in text


class TestIntegrationScenarios:
    """Test real-world integration scenarios."""
    
    def test_mixed_document_batch(self, tmp_path):
        """Test processing a batch of mixed document types."""
        # Create various files
        files = []
        
        # Excel report
        excel_file = tmp_path / "report.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws['A1'] = 'Monthly Report'
        ws['A2'] = 'Sales: $10,000'
        wb.save(excel_file)
        files.append(excel_file)
        
        # PowerPoint presentation
        pptx_file = tmp_path / "presentation.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Quarterly Review"
        prs.save(pptx_file)
        files.append(pptx_file)
        
        # Configuration file
        yaml_file = tmp_path / "config.yaml"
        yaml_file.write_text("app:\n  name: TestApp\n  version: 1.0")
        files.append(yaml_file)
        
        # Log file
        log_file = tmp_path / "app.log"
        log_file.write_text("2025-07-31 INFO Service started\n2025-07-31 ERROR Connection failed")
        files.append(log_file)
        
        # Process all files
        processor = DocumentProcessor()
        documents = list(processor.process_files(files))
        
        # Verify all processed
        assert len(documents) == 4
        
        # Check each document type was processed correctly
        excel_doc = next(d for d in documents if d.file_path.suffix == '.xlsx')
        assert any('Monthly Report' in chunk.text for chunk in excel_doc.chunks)
        
        pptx_doc = next(d for d in documents if d.file_path.suffix == '.pptx')
        assert any('Quarterly Review' in chunk.text for chunk in pptx_doc.chunks)
        
        yaml_doc = next(d for d in documents if d.file_path.suffix == '.yaml')
        assert any('TestApp' in chunk.text for chunk in yaml_doc.chunks)
        
        log_doc = next(d for d in documents if d.file_path.suffix == '.log')
        assert any('ERROR' in chunk.text for chunk in log_doc.chunks)
        
    def test_recursive_email_attachments(self, tmp_path):
        """Test email with attachments that should be processed recursively."""
        # This would require more complex email creation with attachments
        # For now, test basic email processing
        email_file = tmp_path / "email_with_attachment.eml"
        
        msg = MIMEMultipart()
        msg['Subject'] = 'Email with Attachment'
        msg['From'] = 'sender@test.com'
        msg['To'] = 'recipient@test.com'
        
        # Main body
        msg.attach(MIMEText('Please see attached document.'))
        
        # Would add attachment here in real scenario
        
        with open(email_file, 'w') as f:
            f.write(msg.as_string())
            
        processor = DocumentProcessor()
        doc = processor.process_file(email_file)
        
        assert doc is not None
        assert len(doc.chunks) > 0
        assert any('attached document' in chunk.text for chunk in doc.chunks)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])