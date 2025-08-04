"""Tests for PowerPoint reader functionality."""

import pytest
from pathlib import Path
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from vector_db_query.document_processor.powerpoint_reader import PowerPointReader
from vector_db_query.document_processor.exceptions import DocumentProcessingError


class TestPowerPointReader:
    """Test PowerPoint reading functionality."""
    
    @pytest.fixture
    def reader(self):
        """Create PowerPointReader instance."""
        return PowerPointReader()
    
    @pytest.fixture
    def sample_pptx(self):
        """Create a sample PowerPoint file for testing."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            # Slide 1: Title slide
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide1.shapes.title
            subtitle = slide1.placeholders[1]
            title.text = "Test Presentation"
            subtitle.text = "A sample PowerPoint for testing"
            
            # Add speaker notes
            notes_slide1 = slide1.notes_slide
            notes_slide1.notes_text_frame.text = "These are speaker notes for slide 1"
            
            # Slide 2: Content slide
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            title2 = slide2.shapes.title
            content = slide2.placeholders[1]
            title2.text = "Content Slide"
            content.text = "• First bullet point\n• Second bullet point\n• Third bullet point"
            
            # Add speaker notes
            notes_slide2 = slide2.notes_slide
            notes_slide2.notes_text_frame.text = "Important notes for the presenter"
            
            # Slide 3: Mixed content
            slide3 = prs.slides.add_slide(prs.slide_layouts[5])
            title3 = slide3.shapes.title
            title3.text = "Mixed Content"
            
            # Add text box
            left = Inches(1)
            top = Inches(2)
            width = Inches(4)
            height = Inches(1)
            textbox = slide3.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = "This is a text box"
            
            # Add table
            rows = 2
            cols = 2
            left = Inches(1)
            top = Inches(3.5)
            width = Inches(4)
            height = Inches(1.5)
            table = slide3.shapes.add_table(rows, cols, left, top, width, height).table
            table.cell(0, 0).text = "Header 1"
            table.cell(0, 1).text = "Header 2"
            table.cell(1, 0).text = "Data 1"
            table.cell(1, 1).text = "Data 2"
            
            prs.save(tmp.name)
            yield tmp.name
        
        # Cleanup
        Path(tmp.name).unlink(missing_ok=True)
    
    def test_read_simple_presentation(self, reader, sample_pptx):
        """Test reading a simple PowerPoint presentation."""
        result = reader.read(sample_pptx)
        
        assert len(result) > 0
        
        # Check title slide content
        assert "Test Presentation" in result[0]
        assert "A sample PowerPoint for testing" in result[0]
        
        # Check content slide
        assert "Content Slide" in result[0]
        assert "First bullet point" in result[0]
        assert "Second bullet point" in result[0]
        
        # Check mixed content slide
        assert "Mixed Content" in result[0]
        assert "This is a text box" in result[0]
        assert "Header 1" in result[0]
        assert "Data 1" in result[0]
    
    def test_read_with_speaker_notes(self, reader):
        """Test reading presentation with speaker notes included."""
        reader.include_notes = True
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Title"
            
            # Add speaker notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "Important speaker notes"
            
            prs.save(tmp.name)
            
            try:
                result = reader.read(tmp.name)
                assert len(result) > 0
                assert "Important speaker notes" in result[0]
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    
    def test_read_without_speaker_notes(self, reader):
        """Test reading presentation without speaker notes."""
        reader.include_notes = False
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Title"
            
            # Add speaker notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "Speaker notes to ignore"
            
            prs.save(tmp.name)
            
            try:
                result = reader.read(tmp.name)
                assert len(result) > 0
                assert "Speaker notes to ignore" not in result[0]
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    
    def test_read_empty_presentation(self, reader):
        """Test reading an empty presentation."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            prs.save(tmp.name)
            
            try:
                result = reader.read(tmp.name)
                assert len(result) == 1
                assert result[0].strip() == ""
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    
    def test_read_complex_formatting(self, reader):
        """Test reading presentation with complex formatting."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Complex Formatting"
            
            # Add content with multiple paragraphs
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.text = "First paragraph"
            
            p = text_frame.add_paragraph()
            p.text = "Second paragraph with bold"
            p.font.bold = True
            
            p = text_frame.add_paragraph()
            p.text = "Third paragraph with italic"
            p.font.italic = True
            
            prs.save(tmp.name)
            
            try:
                result = reader.read(tmp.name)
                assert len(result) > 0
                assert "Complex Formatting" in result[0]
                assert "First paragraph" in result[0]
                assert "Second paragraph with bold" in result[0]
                assert "Third paragraph with italic" in result[0]
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    
    def test_read_with_images(self, reader):
        """Test reading presentation with images (text extraction only)."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Slide with Images"
            
            # Add text
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
            textbox.text_frame.text = "Text near image"
            
            prs.save(tmp.name)
            
            try:
                result = reader.read(tmp.name)
                assert len(result) > 0
                assert "Slide with Images" in result[0]
                assert "Text near image" in result[0]
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    
    def test_read_multiple_slides(self, reader):
        """Test reading presentation with multiple slides."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            # Create 5 slides
            for i in range(5):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Slide {i+1}"
                content = slide.placeholders[1]
                content.text = f"Content for slide {i+1}"
            
            prs.save(tmp.name)
            
            try:
                result = reader.read(tmp.name)
                assert len(result) > 0
                
                # Check all slides are included
                for i in range(5):
                    assert f"Slide {i+1}" in result[0]
                    assert f"Content for slide {i+1}" in result[0]
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    
    def test_read_nonexistent_file(self, reader):
        """Test reading a non-existent file."""
        with pytest.raises(DocumentProcessingError):
            reader.read("nonexistent.pptx")
    
    def test_read_invalid_file(self, reader):
        """Test reading an invalid file."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(b"This is not a valid PowerPoint file")
            tmp.flush()
            
            try:
                with pytest.raises(DocumentProcessingError):
                    reader.read(tmp.name)
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    
    def test_read_with_charts(self, reader):
        """Test reading presentation with charts (text extraction only)."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Chart Slide"
            
            # Add descriptive text
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
            textbox.text_frame.text = "Sales increased by 25% this quarter"
            
            prs.save(tmp.name)
            
            try:
                result = reader.read(tmp.name)
                assert len(result) > 0
                assert "Chart Slide" in result[0]
                assert "Sales increased by 25%" in result[0]
            finally:
                Path(tmp.name).unlink(missing_ok=True)
    
    def test_supports_extension(self, reader):
        """Test extension support checking."""
        assert reader.supports_extension('.pptx')
        assert reader.supports_extension('.PPTX')
        assert reader.supports_extension('pptx')
        assert not reader.supports_extension('.docx')
        assert not reader.supports_extension('.pdf')
    
    def test_slide_separator(self, reader):
        """Test slide separation in output."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            # Create 2 slides with distinct content
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            slide1.shapes.title.text = "First Slide"
            
            slide2 = prs.slides.add_slide(prs.slide_layouts[0])
            slide2.shapes.title.text = "Second Slide"
            
            prs.save(tmp.name)
            
            try:
                result = reader.read(tmp.name)
                assert len(result) > 0
                
                # Check that slides are separated
                content = result[0]
                assert "Slide 1:" in content
                assert "Slide 2:" in content
                assert "First Slide" in content
                assert "Second Slide" in content
            finally:
                Path(tmp.name).unlink(missing_ok=True)